from pydantic import BaseModel, Field
from typing import Optional, List
from abc import ABC, abstractmethod
import traceback
class PageContent(BaseModel):
    file_name:str = Field(description="文件名")
    page_content:str = Field(description="页面内容")
    page_num:int = Field(description="页码")


class AbsResolver(ABC):
    file_path:list[str] = Field(None, description="文件路径")

    def __init__(self, file_path:list[str]) -> None:
        self.file_path = file_path


    @abstractmethod
    def resolve(self) -> List[PageContent]:
        '''
            解析文件并返回页面内容列表
            file_path: 文件路径
            return: PageContent列表
        '''
        pass


import traceback
import os
from typing import List
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
from pathlib import Path

class DocumentResolver(AbsResolver):
    '''
    简单的根据后缀进行文件类型的判断
    之后，使用对应的解析方法进行文字内容的解析和处理
    '''
    def resolve(self) -> List[PageContent]:
        page_contents = []
        # 路径可能是文件夹，也可能是文件。
        # 如果是文件，直接解析，如果是文件夹，就解析文件夹下的文件
        if os.path.isfile(self.file_path):
            page_contents.extend(self.__resolve_file(file_path=self.file_path))
        else:
            for file_name in os.listdir(self.file_path):
                file_path = os.path.join(self.file_path, file_name)
                page_contents.extend(self.__resolve_file(file_path=file_path))
                
        return page_contents

    def __resolve_file(self, file_path:str):
        page_contents = []
        try:
            file_name = Path(file_path).name
            if file_name.endswith('.txt'):
                with open(file_path, 'r', encoding='utf-8') as file:
                    content = file.read()
                    page_contents.append(PageContent(page_content=content, page_num=1, file_name=file_name))
            elif file_name.endswith('.pdf'):
                doc = fitz.open(file_path)
                # 遍历每一页
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    text = page.get_text()  # 提取文本
                    if text.strip():
                        # 有文本，只取文本数据，图像类的不处理
                        page_contents.append(PageContent(page_content=text, page_num=page_num + 1, file_name=file_name))
                # reader = PdfReader(file_path)
                # for i, page in enumerate(reader.pages):
                #     content = page.extract_text()
                #     page.read_from_stream()
                #     page_contents.append(PageContent(page_content=content, page_num=i + 1, file_name=file_name))
            elif file_name.endswith('.docx'):
                doc = Document(file_path)
                content = '\n'.join([para.text for para in doc.paragraphs])
                page_contents.append(PageContent(page_content=content, page_num=1, file_name=file_name))
            elif file_name.endswith('.pptx'):
                presentation = Presentation(file_path)
                for i, slide in enumerate(presentation.slides):
                    content = '\n'.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                    page_contents.append(PageContent(page_content=content, page_num=i + 1, file_name=file_name))
                            
            # 可以根据需要添加更多文件类型的解析
                        
        except Exception as e:
            print(f"文件读取出错==,异常信息：{traceback.format_exc()}")

        return page_contents

if "__main__" == __name__:
    result:list[PageContent] = DocumentResolver(file_path="/Users/zhanhd/Downloads").resolve()

    for page in result:
        print(page.model_dump())